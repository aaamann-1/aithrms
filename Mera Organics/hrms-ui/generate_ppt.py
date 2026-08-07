from pptx import Presentation
from pptx.util import Pt
from pathlib import Path

output = Path(r'C:\Users\asus\OneDrive\Desktop\Mera Organics\hrms-ui\HRMS_Portal_Demo.pptx')
prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'HRMS Portal Demo'
slide.placeholders[1].text = 'Static UI prototype with role-based portals and demo auth flow.'


def add_slide(title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(18)


add_slide(
    'Project Overview',
    [
        'Browser-based HRMS portal prototype',
        'No backend yet: demo auth using localStorage/sessionStorage',
        'Role-aware portal for Admin, Team Leader, Employee',
        'Landing page, login, signup, portal dashboard',
    ],
)

add_slide(
    'Key Pages',
    [
        'Landing page: product intro and navigation',
        'Login page: demo credentials + redirect guard',
        'Signup page: employee self-registration',
        'Portal page: role-based navigation and actions',
    ],
)

add_slide(
    'Authentication Flow',
    [
        'Users stored in browser localStorage for demo',
        'SessionStorage holds current logged-in user',
        'Redirects protect portal access from unauthenticated users',
        'Admin can create new employee/team leader logins',
    ],
)

add_slide(
    'Roles & Permissions',
    [
        'Admin: full dashboard, reports, settings, employee management',
        'Team Leader: team employee access, issues view, no leave/attendance approval',
        'Employee: self profile, mark attendance, apply leave, report issues',
    ],
)

add_slide(
    'Portal Features',
    [
        'Dashboard: summaries, announcements, stats',
        'Employees: directory and profile view',
        'Attendance: overview for managers, self-mark for employees',
        'Leave: admin approval, employee requests, team leader view only',
        'Issues & Reports: role-specific task and report panels',
    ],
)

add_slide(
    'Demo Credentials & Next Steps',
    [
        'Admin: admin@hrms.local / admin123',
        'Team Leader: teamleader@hrms.local / leader123',
        'Employee: employee@hrms.local / employee123',
        'Next: add real backend, role-based data, persistence, UI polish',
    ],
)

prs.save(output)
print(f'Saved presentation to {output}')
